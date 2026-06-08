"""Fill a GEP template into a branded output deck.

Spec JSON:
{
  "template": "_assets/Assets/3Pt.pptx",
  "output":   "out.pptx",
  "slides": [
    {"index": 1,
     "text":   {"2": "My Title", "105": "Discover", "75": "body copy..."},
     "images": {"49": "C:/path/icon.png"}}
  ]
}

- Keeps ONLY the listed slides, in the listed order.
- text: maps shape_id -> new string. Replaces text while preserving the
  formatting of the shape's first run (font, size, color, bold).
- images: maps picture shape_id -> image path. Swaps the picture's blip.
Usage: python fill_slide.py spec.json
"""
import sys, json, copy
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

def find_shape(shapes, sid):
    for sh in shapes:
        if sh.shape_id == sid:
            return sh
        if sh.shape_type == 6:  # group
            r = find_shape(sh.shapes, sid)
            if r is not None:
                return r
    return None

def _set_para(p, value):
    """Set a paragraph's text to `value`, preserving its first run's formatting."""
    if p.runs:
        p.runs[0].text = value
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = value

def set_text(shape, value):
    """Set shape text. `value` is a str (whole frame -> one paragraph) or a
    list of strings (one per paragraph, preserving each paragraph's format)."""
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    if isinstance(value, list):
        paras = tf.paragraphs
        n = len(value)
        # drop paragraphs beyond what we need
        for p in paras[n:]:
            p._p.getparent().remove(p._p)
        for i, val in enumerate(value):
            if i < len(tf.paragraphs):
                _set_para(tf.paragraphs[i], val)
            else:
                _set_para(tf.add_paragraph(), val)
        return
    p0 = tf.paragraphs[0]
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    _set_para(p0, value)

def swap_image(prs, shape, img_path):
    # replace the image part's blob by reusing the existing relationship
    rId = shape._element.blipFill.blip.embed
    img_part, _ = shape.part.get_or_add_image_part(img_path)
    # rebind relationship target
    shape.part.rels._rels[rId]._target = img_part
    shape._element.blipFill.blip.embed = shape.part.relate_to(
        img_part, shape.part.rels._rels[rId].reltype)

def keep_slides(prs, indices_1based):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    keep = [ids[i-1] for i in indices_1based]
    for sid in ids:
        sldIdLst.remove(sid)
    for sid in keep:
        sldIdLst.append(sid)

def main():
    spec = json.load(open(sys.argv[1], encoding="utf-8"))
    prs = Presentation(spec["template"])
    slides = list(prs.slides)
    for s in spec["slides"]:
        slide = slides[s["index"]-1]
        for sid, val in s.get("text", {}).items():
            sh = find_shape(slide.shapes, int(sid))
            if sh is None or not sh.has_text_frame:
                print(f"WARN slide {s['index']} text shape {sid} not found/no text"); continue
            set_text(sh, val)
        for sid, path in s.get("images", {}).items():
            sh = find_shape(slide.shapes, int(sid))
            if sh is None:
                print(f"WARN slide {s['index']} image shape {sid} not found"); continue
            swap_image(prs, sh, path)
        for img in s.get("add_images", []):
            slide.shapes.add_picture(
                img["path"], Inches(img["L"]), Inches(img["T"]),
                Inches(img["W"]), Inches(img["H"]))
        for sid in s.get("remove_shapes", []):
            sh = find_shape(slide.shapes, int(sid))
            if sh is None:
                print(f"WARN slide {s['index']} remove shape {sid} not found"); continue
            sh._element.getparent().remove(sh._element)
    keep_slides(prs, [s["index"] for s in spec["slides"]])
    prs.save(spec["output"])
    print("SAVED", spec["output"])

if __name__ == "__main__":
    main()
