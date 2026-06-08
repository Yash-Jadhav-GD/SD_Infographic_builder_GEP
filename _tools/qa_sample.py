"""Fill a representative SAMPLE of auto-mapped templates, one slide per file.

Uses the production path (fill + keep only that slide) which PowerPoint can open.
Picks up to PER_BUCKET entries for each item_count so QA spans the variety.
Usage: python qa_sample.py <deck> <catalog_auto.json> <out_dir> [per_bucket]
"""
import sys, json, copy
from pathlib import Path
from collections import defaultdict
from pptx import Presentation
from fill_slide import find_shape, set_text, keep_slides

deck, cat_path, out_dir = sys.argv[1], sys.argv[2], sys.argv[3]
PER = int(sys.argv[4]) if len(sys.argv) > 4 else 3
out_dir = Path(out_dir); out_dir.mkdir(parents=True, exist_ok=True)
cat = json.load(open(cat_path, encoding="utf-8"))["templates"]

buckets = defaultdict(list)
for e in cat:
    buckets[e["item_count"]].append(e)
sample = []
for ic in sorted(buckets):
    sample += buckets[ic][:PER]

manifest = []
for e in sample:
    prs = Presentation(deck)
    slide = list(prs.slides)[e["slide_index"] - 1]
    if "title_shape" in e:
        sh = find_shape(slide.shapes, int(e["title_shape"]))
        if sh is not None and sh.has_text_frame:
            set_text(sh, e["name"])
    for i, slot in enumerate(e["items"], start=1):
        label = f"LABEL {i}"
        body = f"Body {i}: sample copy confirming slot {i} maps correctly here."
        if "combo_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["combo_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, [label, body])
        if "name_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["name_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, label)
        if "body_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["body_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, body)
    keep_slides(prs, [e["slide_index"]])
    fn = out_dir / f"s{e['slide_index']:03d}_{e['id'][:30]}.pptx"
    prs.save(str(fn))
    manifest.append({"file": fn.name, "slide": e["slide_index"],
                     "name": e["name"], "items": e["item_count"]})

json.dump(manifest, open(out_dir / "_manifest.json", "w", encoding="utf-8"), indent=2)
print(f"wrote {len(manifest)} sample decks to {out_dir}")
