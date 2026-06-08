"""Auto-generate catalog entries for a clean template deck.

The 'Infographic Templates' deck is highly regular:
  - top-left header text  = the layout NAME (e.g. "Doughnut 1")
  - "Lorem ipsum"         = an item label/name slot
  - "Lorem ipsum dolor …" = an item body slot
  - "N. Lorem ipsum"      = a numbered item label

Heuristic per slide:
  title_shape  -> the header-name text shape (top band, not a lorem placeholder)
  name slots   -> short lorem-label text shapes
  body slots   -> lorem-body text shapes
  items        -> names paired to nearest body, ordered in reading order

Output: backend/catalog_auto.json  (+ a printed summary).
Usage:  python auto_catalog.py <deck.pptx> <template_rel_path> <out.json>
"""
import sys, json, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

NAME_RE = re.compile(r"^\s*(?:\d+[\.\)]\s*)?lorem ipsum\s*$", re.I)
BODY_RE = re.compile(r"^\s*lorem ipsum dolor", re.I)


def emu_in(v):
    return round(Emu(v).inches, 3) if v is not None else 0.0


def walk(shapes, depth=0):
    out = []
    for sh in shapes:
        txt = ""
        if sh.has_text_frame:
            txt = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text).strip()
        is_pic = sh.shape_type == 13
        out.append({
            "id": sh.shape_id, "L": emu_in(sh.left), "T": emu_in(sh.top),
            "W": emu_in(sh.width), "H": emu_in(sh.height),
            "pic": is_pic, "text": txt,
            "cx": emu_in(sh.left) + emu_in(sh.width) / 2,
            "cy": emu_in(sh.top) + emu_in(sh.height) / 2,
        })
        if sh.shape_type == 6:  # group -> recurse
            out += walk(sh.shapes, depth + 1)
    return out


def slug(s):
    s = re.sub(r"[^a-z0-9]+", "-", s.lower()).strip("-")
    return s or "layout"


def classify(rows):
    title, names, bodies, combos = None, [], [], []
    for r in rows:
        t = r["text"]
        if not t:
            continue
        low = t.lower()
        lines = [ln for ln in t.splitlines() if ln.strip()]
        if "lorem ipsum dolor" in low:
            if low.lstrip().startswith("lorem ipsum dolor"):
                bodies.append(r)               # body only
            else:
                # label in first line, body in the rest -> combo shape
                r["_label"] = lines[0]
                r["_body"] = " ".join(lines[1:])
                combos.append(r)
        elif NAME_RE.match(t):
            names.append(r)                    # label only
        elif r["T"] < 0.95 and "lorem" not in low and len(t) < 60:
            if title is None or r["T"] < title["T"]:
                title = r                       # top-band header = title slot
    return title, names, bodies, combos


def build_entry(idx, rows, template_rel):
    title, names, bodies, combos = classify(rows)
    if not bodies and not names and not combos:
        return None, "no item placeholders"
    used_bodies = set()
    items = []
    for c in combos:
        items.append({"combo_shape": c["id"], "_t": c["T"], "_l": c["L"],
                      "_nt": c.get("_label", ""), "_bt": c.get("_body", "")})
    # pair each name with nearest unused body
    for nm in names:
        best, bestd = None, 1e9
        for i, b in enumerate(bodies):
            if i in used_bodies:
                continue
            d = (nm["cx"] - b["cx"]) ** 2 + (nm["cy"] - b["cy"]) ** 2
            if d < bestd:
                bestd, best = d, i
        slot = {"name_shape": nm["id"], "_t": nm["T"], "_l": nm["L"],
                "_nt": nm["text"]}
        if best is not None and bestd < 16:  # within ~4in
            used_bodies.add(best)
            slot["body_shape"] = bodies[best]["id"]
            slot["_bt"] = bodies[best]["text"]
        items.append(slot)
    # leftover bodies become body-only items
    for i, b in enumerate(bodies):
        if i not in used_bodies:
            items.append({"body_shape": b["id"], "_t": b["T"], "_l": b["L"],
                          "_bt": b["text"]})
    # reading order: row bands (0.5in) then left-to-right
    items.sort(key=lambda s: (round(s["_t"] / 0.5), s["_l"]))

    name_cap = body_cap = 0
    clean = []
    for s in items:
        c = {}
        if "combo_shape" in s:
            c["combo_shape"] = s["combo_shape"]
            name_cap = max(name_cap, int(len(s.get("_nt", "")) * 1.5), 24)
            body_cap = max(body_cap, int(len(s.get("_bt", "")) * 1.3), 60)
        if "name_shape" in s:
            c["name_shape"] = s["name_shape"]
            name_cap = max(name_cap, int(len(s.get("_nt", "")) * 1.5), 24)
        if "body_shape" in s:
            c["body_shape"] = s["body_shape"]
            body_cap = max(body_cap, int(len(s.get("_bt", "")) * 1.3), 60)
        clean.append(c)

    layout_name = (title["text"].strip() if title else f"Template {idx}")
    entry = {
        "id": f"{slug(layout_name)}-s{idx}",
        "name": layout_name,
        "family": "auto",
        "item_count": len(clean),
        "template": template_rel,
        "slide_index": idx,
        "description": f"Auto-mapped layout '{layout_name}' ({len(clean)} items).",
        "items": clean,
        "name_capacity": name_cap or 24,
        "body_capacity": body_cap or 120,
        "auto": True,
    }
    if title:
        entry["title_shape"] = title["id"]
    return entry, None


def main():
    deck, template_rel, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    prs = Presentation(deck)
    entries, skipped = [], []
    for i, slide in enumerate(prs.slides, start=1):
        rows = walk(slide.shapes)
        entry, why = build_entry(i, rows, template_rel)
        if entry is None:
            skipped.append((i, why))
            continue
        entries.append(entry)
    json.dump({"templates": entries}, open(out_path, "w", encoding="utf-8"), indent=2)
    # summary
    from collections import Counter
    counts = Counter(e["item_count"] for e in entries)
    print(f"slides={len(prs.slides)} mapped={len(entries)} skipped={len(skipped)}")
    print("item_count distribution:", dict(sorted(counts.items())))
    print("with_title:", sum(1 for e in entries if "title_shape" in e))
    print("sample names:", [e["name"] for e in entries[:12]])
    if skipped:
        print("skipped slides:", skipped[:20])


if __name__ == "__main__":
    main()
