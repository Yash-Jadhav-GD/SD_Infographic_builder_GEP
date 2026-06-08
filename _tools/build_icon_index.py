"""Extract every icon from Icons.pptx and tag it with keywords.

Each slide = a background keyword Table (one keyword blob per row) + a grid of
icon Pictures laid over it. We map each picture to the table row that vertically
contains it, and inherit that row's keywords as the icon's tags.

Outputs:
  <out_dir>/icons/<slide>_<shapeid>.<ext>   extracted icon images
  <out_dir>/icons_index.json                [{file, slide, keywords, tags}]

Usage: python build_icon_index.py <Icons.pptx> <out_dir>
"""
import sys, os, json, re
from pptx import Presentation
from pptx.util import Emu

STOP = set("the a an of and or to for with in on at by".split())

def tokenize(text):
    toks = re.findall(r"[A-Za-z][A-Za-z0-9]+", (text or "").lower())
    seen, out = set(), []
    for t in toks:
        if t in STOP or len(t) < 2 or t in seen:
            continue
        seen.add(t); out.append(t)
    return out

def table_rows(tbl_frame):
    """Return list of (y_top_in, y_bot_in, keyword_text) per table row."""
    tbl = tbl_frame.table
    top = Emu(tbl_frame.top).inches
    rows = []
    y = top
    n = len(tbl.rows)
    total_h = Emu(tbl_frame.height).inches
    for i, row in enumerate(tbl.rows):
        h = Emu(row.height).inches if row.height else (total_h / n)
        text = " ".join(c.text for c in row.cells if c.text).strip()
        rows.append((y, y + h, text))
        y += h
    return rows

def main():
    src, out_dir = sys.argv[1], sys.argv[2]
    icons_dir = os.path.join(out_dir, "icons")
    os.makedirs(icons_dir, exist_ok=True)
    prs = Presentation(src)
    index = []
    for si, slide in enumerate(prs.slides, start=1):
        rows = []
        for sh in slide.shapes:
            if sh.has_table:
                rows = table_rows(sh)
                break
        for sh in slide.shapes:
            if sh.shape_type != 13:  # picture
                continue
            try:
                img = sh.image
            except Exception:
                continue
            cy = Emu(sh.top).inches + Emu(sh.height).inches / 2
            kw = ""
            for (yt, yb, text) in rows:
                if yt <= cy <= yb and text:
                    kw = text; break
            if not kw and rows:  # nearest row fallback
                kw = min(rows, key=lambda r: abs((r[0]+r[1])/2 - cy))[2]
            ext = img.ext or "png"
            fname = f"{si}_{sh.shape_id}.{ext}"
            with open(os.path.join(icons_dir, fname), "wb") as f:
                f.write(img.blob)
            index.append({"file": f"icons/{fname}", "slide": si,
                          "keywords": kw, "tags": tokenize(kw)})
    with open(os.path.join(out_dir, "icons_index.json"), "w", encoding="utf-8") as f:
        json.dump(index, f, indent=1)
    tagged = sum(1 for e in index if e["tags"])
    print(f"icons extracted: {len(index)} | with tags: {tagged}")

if __name__ == "__main__":
    main()
