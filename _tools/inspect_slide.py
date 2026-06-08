"""Inspect a template slide: list every text/picture shape with position + current text.
Usage: python inspect_slide.py <deck.pptx> <slide_index_1based>
Positions are in inches. Use the printed shape_id to target shapes when filling.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

def emu_in(v):
    return round(Emu(v).inches, 2) if v is not None else None

def walk(shapes, depth=0):
    rows = []
    for sh in shapes:
        kind = sh.shape_type
        txt = ""
        if sh.has_text_frame:
            txt = " | ".join(p.text for p in sh.text_frame.paragraphs if p.text).strip()
        is_pic = (str(kind) == "PICTURE (13)") or sh.shape_type == 13
        rows.append({
            "id": sh.shape_id, "name": sh.name, "type": str(kind),
            "L": emu_in(sh.left), "T": emu_in(sh.top),
            "W": emu_in(sh.width), "H": emu_in(sh.height),
            "pic": is_pic, "text": txt, "depth": depth,
        })
        if sh.shape_type == 6:  # group
            rows += walk(sh.shapes, depth+1)
    return rows

def main():
    deck, idx = sys.argv[1], int(sys.argv[2])
    prs = Presentation(deck)
    slide = prs.slides[idx-1]
    rows = walk(slide.shapes)
    print(f"# {deck} slide {idx} — {len(rows)} shapes")
    for r in rows:
        pad = "  "*r["depth"]
        flag = "[PIC]" if r["pic"] else ("[TXT]" if r["text"] else "     ")
        t = (r["text"][:70] + "…") if len(r["text"])>70 else r["text"]
        print(f'{r["id"]:>4} {flag} {pad}{r["name"][:24]:<24} @({r["L"]},{r["T"]}) {r["W"]}x{r["H"]}  {t}')

if __name__ == "__main__":
    main()
