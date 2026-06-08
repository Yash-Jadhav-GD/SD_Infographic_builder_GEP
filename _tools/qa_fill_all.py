"""Fill EVERY auto-mapped template slide with dummy data, in one deck copy.

Lets us render the whole deck and eyeball which slides mapped cleanly.
Distinctive dummy text makes mis-maps obvious (wrong slot = wrong text).
Usage: python qa_fill_all.py <deck.pptx> <catalog_auto.json> <out.pptx>
"""
import sys, json
from pptx import Presentation
from fill_slide import find_shape, set_text

deck, cat_path, out = sys.argv[1], sys.argv[2], sys.argv[3]
LIMIT = int(sys.argv[4]) if len(sys.argv) > 4 else 0   # only fill entries whose slide_index <= LIMIT (0 = all)
cat = json.load(open(cat_path, encoding="utf-8"))["templates"]
if LIMIT:
    cat = [e for e in cat if e["slide_index"] <= LIMIT]
prs = Presentation(deck)
slides = list(prs.slides)

for e in cat:
    slide = slides[e["slide_index"] - 1]
    if "title_shape" in e:
        sh = find_shape(slide.shapes, int(e["title_shape"]))
        if sh is not None and sh.has_text_frame:
            set_text(sh, e["name"])  # keep layout name as the title
    for i, slot in enumerate(e["items"], start=1):
        label = f"LABEL {i}"
        body = f"Body {i}: sample copy to confirm this slot maps to item {i}."
        if "combo_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["combo_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, [label, body])
        if "name_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["name_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, label)
        if "body_shape" in slot:
            sh = find_shape(slide.shapes, int(slot["body_shape"]))
            if sh is not None and sh.has_text_frame:
                set_text(sh, body)

prs.save(out)
print("SAVED", out, "filled", len(cat), "slides")
